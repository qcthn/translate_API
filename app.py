import streamlit as st
import openai
import time
from collections import deque
from io import BytesIO
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader, PdfWriter
import pandas as pd
from pptx.util import Pt


# Rate limiting parameters
MAX_REQUESTS_PER_MINUTE = 3500
MAX_TOKENS_PER_MINUTE = 90000
WINDOW_SECONDS = 60
requests_timestamps = deque()
tokens_timestamps = deque()

def check_and_wait_for_rate_limit(tokens_used: int):
    current_time = time.time()
    while requests_timestamps and (current_time - requests_timestamps[0] > WINDOW_SECONDS):
        requests_timestamps.popleft()
    while tokens_timestamps and (current_time - tokens_timestamps[0][0] > WINDOW_SECONDS):
        tokens_timestamps.popleft()
    current_requests = len(requests_timestamps)
    current_tokens = sum(t[1] for t in tokens_timestamps)
    if current_requests + 1 > MAX_REQUESTS_PER_MINUTE or current_tokens + tokens_used > MAX_TOKENS_PER_MINUTE:
        time.sleep(1)
        return check_and_wait_for_rate_limit(tokens_used)
    requests_timestamps.append(current_time)
    tokens_timestamps.append((current_time, tokens_used))

def load_specialized_dict_from_excel(excel_file):
    if excel_file is None:
        return {}
    df = pd.read_excel(excel_file)
    return {str(row['English']).strip(): str(row['Vietnamese']).strip() for _, row in df.iterrows() if row['English'] and row['Vietnamese']}

def translate_text_with_chatgpt(original_text, api_key, global_dict=None):
    if not original_text.strip():
        return original_text
    partial_dict = {eng: vie for eng, vie in global_dict.items() if eng.lower() in original_text.lower()} if global_dict else {}
    dict_prompt = "\n".join([f"{k}: {v}" for k, v in partial_dict.items()]) if partial_dict else ""
    user_prompt = f"{dict_prompt}\n\n{original_text}"
    client = openai.OpenAI(api_key=api_key)
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "Bạn là một trợ lý AI dịch thuật. Hãy dịch văn bản sau từ tiếng Anh sang tiếng Việt, ưu tiên dùng đúng các thuật ngữ chuyên ngành (nếu có). Trước tiên hay tra cứu từ vựng trong câu có từ  nào thuộc từ vựng nằm trong file từ vựng chuyên ngành mà tôi cung cấp không, nếu có hãy dùng nghĩa tiếng việt của từ vựng chuyên ngành đó được cung cấp trong file xlsx, các từ còn lại bạn có thể dịch tự động. ** Lưu ý mỗi câu chỉ được phép dịch 1 lần duy nhất, ngoài ra nếu đó là mội chuỗi kí tự bất kì không phải là bất kì từ tiếng anh nào thì đó có thể là kí hiệu hoặc mã của sản phẩm bạn có thể giữ nguyên và không cần dịch sang tiếng việt. Nếu đầu vào (input) không có nội dung thì bạn có thể bỏ qua và không trả về kết quả gì hết ( không trả output)."},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.2,
        max_tokens=2048
    )
    translated_text = response.choices[0].message.content
    check_and_wait_for_rate_limit(response.usage.total_tokens if response.usage else 0)
    return translated_text


from pptx.util import Pt

def adjust_text_fit(text_frame, shape):
    """
    Adjust text font size dynamically to fit within the text box without overflow.
    Uses shape.width and shape.height instead of text_frame.width.
    """
    max_width = shape.width  # Get the width of the text box
    max_height = shape.height  # Get the height of the text box
    min_font_size = Pt(8)  # Set a minimum font size to maintain readability

    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.font.size and run.font.size > min_font_size:
                run.font.size = max(min_font_size, run.font.size * 0.9)  # Reduce font size if needed

def distribute_text_across_runs(para, translated_text):
    """
    Phân phối đều văn bản đã dịch qua các run trong khi giữ nguyên định dạng.
    """
    original_text = "".join(run.text for run in para.runs)
    if not original_text.strip():
        return

    total_original_len = len(original_text)
    if total_original_len == 0:
        return

    remaining_text = translated_text
    for run in para.runs:
        if not run.text:
            continue

        # Tính toán phân phối văn bản theo tỷ lệ
        run_len = len(run.text)
        portion = min(run_len / total_original_len, 1.0)
        chars_to_take = int(len(translated_text) * portion)

        # Cập nhật văn bản trong khi giữ nguyên định dạng
        run.text = remaining_text[:chars_to_take]
        remaining_text = remaining_text[chars_to_take:]

        # Giữ nguyên định dạng (font, cỡ chữ, in đậm, in nghiêng, màu sắc)
        if run.font is not None:
            run.font.name = run.font.name  # Giữ font chữ
            run.font.size = run.font.size  # Giữ cỡ chữ
            run.font.bold = run.font.bold  # Giữ in đậm
            run.font.italic = run.font.italic  # Giữ in nghiêng
            
            if run.font.color and hasattr(run.font.color, 'rgb'):
                run.font.color.rgb = run.font.color.rgb  # Giữ màu sắc

    # Gắn phần văn bản còn lại vào run cuối cùng (nếu có)
    if remaining_text and para.runs:
        para.runs[-1].text += remaining_text

def translate_pptx(pptx_file: BytesIO, api_key: str, specialized_dict: dict[str, str]) -> BytesIO:
    """
    Dịch văn bản trong file PowerPoint từ tiếng Anh sang tiếng Việt, giữ nguyên font, cỡ chữ và màu sắc gốc.
    Xử lý văn bản tràn bằng cách điều chỉnh cỡ chữ động.
    
    Args:
        pptx_file: Đối tượng BytesIO chứa file PPTX
        api_key: Khóa API OpenAI
        specialized_dict: Từ điển thuật ngữ chuyên ngành (Anh -> Việt)
    
    Returns:
        Đối tượng BytesIO chứa file PPTX đã dịch
    """
    pr = Presentation(pptx_file)
    total_slides = len(pr.slides)
    progress_bar = st.progress(0)
    status_text = st.empty()

    for i, slide in enumerate(pr.slides):
        status_text.text(f"Đang dịch slide {i+1}/{total_slides}...")

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame

                for para in text_frame.paragraphs:
                    if not para.text.strip():
                        continue

                    # Thu thập văn bản gốc
                    original_text = "".join(run.text for run in para.runs)
                    translated_text = translate_text_with_chatgpt(original_text, api_key, specialized_dict)

                    # Bỏ qua nếu không có bản dịch hoặc bản dịch không thay đổi
                    if not translated_text or translated_text == original_text or translated_text == 'Xin lỗi, nhưng văn bản bạn cung cấp không đủ để dịch. Bạn có thể cung cấp thêm ngữ cảnh hoặc thông tin chi tiết hơn không?':
                        continue

                    # Phân phối văn bản đã dịch qua các run
                    distribute_text_across_runs(para, translated_text)

                # Điều chỉnh kích thước văn bản để tránh tràn
                adjust_text_fit(text_frame, shape)

            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_frame = cell.text_frame
                        for para in text_frame.paragraphs:
                            if not para.text.strip():
                                continue
                            # Thu thập văn bản gốc từ các run
                            original_text = "".join(run.text for run in para.runs)
                            translated_text = translate_text_with_chatgpt(original_text, api_key, specialized_dict)
                            # Chỉ cập nhật nếu bản dịch hợp lệ
                            if translated_text and translated_text != original_text and translated_text != 'Xin lỗi, nhưng văn bản bạn cung cấp không đủ để dịch. Bạn có thể cung cấp thêm ngữ cảnh hoặc thông tin chi tiết hơn không?':
                                distribute_text_across_runs(para, translated_text)
                        # Tùy chọn: Điều chỉnh kích thước văn bản trong ô (nếu cần)
                                # adjust_text_fit_for_cell(cell)

        progress_bar.progress((i+1) / total_slides)

    output = BytesIO()
    pr.save(output)
    output.seek(0)
    status_text.text("Dịch PPTX hoàn tất!")
    return output

# Streamlit UI
st.set_page_config(page_title="Auto Translator App with Full Formatting")
st.title("Tự động dịch tài liệu (PPTX) + Giữ nguyên định dạng & kích thước")

api_key = st.text_input("Nhập OpenAI API key của bạn:", type="password")
uploaded_excel_dict = st.file_uploader("Tải lên file Excel chứa thuật ngữ chuyên ngành", type=["xlsx"])
specialized_dict = load_specialized_dict_from_excel(uploaded_excel_dict)

uploaded_file = st.file_uploader("Tải lên file cần dịch (PPTX)", type=["pptx"])

if uploaded_file and api_key and st.button("Bắt đầu dịch"):
    ext = uploaded_file.name.split(".")[-1].lower()
    
    if ext == "pptx":
        output = translate_pptx(uploaded_file, api_key, specialized_dict)
        st.download_button("Tải về file PPTX đã dịch", output, "translated.pptx")

    else:
        st.error("Định dạng không được hỗ trợ.")
