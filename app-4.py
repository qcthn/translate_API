import streamlit as st
import openai
import os
import time
from collections import deque
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from PyPDF2 import PdfReader, PdfWriter
import pandas as pd

# Rate limiting parameters
MAX_REQUESTS_PER_MINUTE = 3500
MAX_TOKENS_PER_MINUTE = 90000
WINDOW_SECONDS = 60
requests_timestamps = deque()
tokens_timestamps = deque()

def check_and_wait_for_rate_limit(tokens_used: int):
    current_time = time.time()
    while requests_timestamps and (current_time - requests_timestamps[0] > WINDOW_SECONDS):
        requests_timestamps.popleft()
    while tokens_timestamps and (current_time - tokens_timestamps[0][0] > WINDOW_SECONDS):
        tokens_timestamps.popleft()
    current_requests = len(requests_timestamps)
    current_tokens = sum(t[1] for t in tokens_timestamps)
    if current_requests + 1 > MAX_REQUESTS_PER_MINUTE or current_tokens + tokens_used > MAX_TOKENS_PER_MINUTE:
        time.sleep(1)
        return check_and_wait_for_rate_limit(tokens_used)
    requests_timestamps.append(current_time)
    tokens_timestamps.append((current_time, tokens_used))

def load_specialized_dict_from_excel(excel_file):
    if excel_file is None:
        return {}
    df = pd.read_excel(excel_file)
    return {str(row['English']).strip(): str(row['Vietnamese']).strip() for _, row in df.iterrows() if row['English'] and row['Vietnamese']}

def translate_text_with_chatgpt(original_text, api_key, global_dict=None):
    if not original_text.strip():
        return original_text
    partial_dict = {eng: vie for eng, vie in global_dict.items() if eng.lower() in original_text.lower()} if global_dict else {}
    dict_prompt = "\n".join([f"{k}: {v}" for k, v in partial_dict.items()]) if partial_dict else ""
    user_prompt = f"{dict_prompt}\n\n{original_text}"
    client = openai.OpenAI(api_key=api_key)
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "Bạn là một trợ lý AI dịch thuật. Hãy dịch văn bản sau từ tiếng Anh sang tiếng Việt, ưu tiên dùng đúng các thuật ngữ chuyên ngành (nếu có). Trước tiên hay tra cứu từ vựng trong câu có từ  nào thuộc từ vựng nằm trong file từ vựng chuyên ngành mà tôi cung cấp không, nếu có hãy dùng nghĩa tiếng việt của từ vựng chuyên ngành đó được cung cấp trong file xlsx, các từ còn lại bạn có thể dịch tự động. ** Lưu ý mỗi câu chỉ được phép dịch 1 lần duy nhất, ngoài ra nếu đó là mội chuỗi kí tự bất kì không phải là bất kì từ tiếng anh nào thì đó có thể là kí hiệu hoặc mã của sản phẩm bạn có thể giữ nguyên và không cần dịch sang tiếng việt. Nếu đầu vào (input) không có nội dung thì bạn có thể bỏ qua và không trả về kết quả gì hết ( không trả output)"},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.2,
        max_tokens=2048
    )
    translated_text = response.choices[0].message.content
    check_and_wait_for_rate_limit(response.usage.total_tokens if response.usage else 0)
    return translated_text

def translate_pptx(pptx_file: BytesIO, api_key: str, specialized_dict: dict[str, str]) -> BytesIO:
    """
    Translate text in a PowerPoint file while preserving original font, size, and color.
    
    Args:
        pptx_file: BytesIO object containing the PPTX file
        api_key: OpenAI API key
        specialized_dict: Dictionary of specialized terms (English -> Vietnamese)
    
    Returns:
        BytesIO object with the translated PPTX file
    """
    pr = Presentation(pptx_file)

    for slide in pr.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame
            for para in text_frame.paragraphs:
                if not para.text.strip():
                    continue

                # Collect all runs' text while preserving formatting
                original_text = "".join(run.text for run in para.runs)
                translated_text = translate_text_with_chatgpt(original_text, api_key, specialized_dict)

                # Skip if translation is empty or unchanged
                if not translated_text or translated_text == original_text:
                    continue

                # Distribute translated text across runs proportionally
                total_original_len = len(original_text)
                if total_original_len == 0:
                    continue

                remaining_text = translated_text
                for run in para.runs:
                    if not run.text:
                        continue

                    # Calculate portion of translated text for this run
                    run_len = len(run.text)
                    portion = min(run_len / total_original_len, 1.0)
                    chars_to_take = int(len(translated_text) * portion)

                    # Update run text while preserving formatting
                    run.text = remaining_text[:chars_to_take]
                    remaining_text = remaining_text[chars_to_take:]

                # If there's remaining text, append it to the last run
                if remaining_text and para.runs:
                    para.runs[-1].text += remaining_text

    # Save to BytesIO
    output = BytesIO()
    pr.save(output)
    output.seek(0)
    return output

def translate_docx(docx_file, api_key, specialized_dict):
    doc = Document(docx_file)
    for para in doc.paragraphs:
        if para.text.strip():
            translated_text = translate_text_with_chatgpt(para.text, api_key, specialized_dict)
            para.text = translated_text
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output

def translate_pdf(pdf_file, api_key, specialized_dict):
    reader = PdfReader(pdf_file)
    writer = PdfWriter()
    for page in reader.pages:
        text = page.extract_text()
        if text:
            translated_text = translate_text_with_chatgpt(text, api_key, specialized_dict)
            writer.add_page(page)
            writer.pages[-1].merge_text(translated_text)
    output = BytesIO()
    writer.write(output)
    output.seek(0)
    return output

st.set_page_config(page_title="Auto Translator App with Format Preservation")
st.title("Tự động dịch tài liệu (PPTX) + Giữ định dạng")
api_key = st.text_input("Nhập OpenAI API key của bạn:", type="password")
uploaded_excel_dict = st.file_uploader("Tải lên file Excel chứa thuật ngữ chuyên ngành", type=["xlsx"])
specialized_dict = load_specialized_dict_from_excel(uploaded_excel_dict)
uploaded_file = st.file_uploader("Tải lên file cần dịch (PPTX)", type=["pptx", "docx", "pdf"])
if uploaded_file and api_key and st.button("Bắt đầu dịch"):
    ext = uploaded_file.name.split(".")[-1].lower()
    if ext == "pptx":
        output = translate_pptx(uploaded_file, api_key, specialized_dict)
        st.download_button("Tải về file PPTX đã dịch", output, "translated.pptx")
    elif ext == "docx":
        output = translate_docx(uploaded_file, api_key, specialized_dict)
        st.download_button("Tải về file DOCX đã dịch", output, "translated.docx")
    elif ext == "pdf":
        output = translate_pdf(uploaded_file, api_key, specialized_dict)
        st.download_button("Tải về file PDF đã dịch", output, "translated.pdf")
    else:
        st.error("Định dạng không được hỗ trợ.")
