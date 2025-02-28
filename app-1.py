import streamlit as st
import openai
import os
import time
from collections import deque
from io import BytesIO

# Thư viện cho Word, PowerPoint, PDF
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader, PdfWriter

# Đọc file Excel
import pandas as pd

########################################################################
# 1) CẤU HÌNH GIỚI HẠN (RATE LIMIT) THEO YÊU CẦU
########################################################################
MAX_REQUESTS_PER_MINUTE = 3500  # GPT-3.5-turbo RPM
MAX_TOKENS_PER_MINUTE = 90000   # GPT-3.5-turbo TPM
WINDOW_SECONDS = 60             # Thời gian cửa sổ để tính giới hạn (60 giây)

requests_timestamps = deque()
tokens_timestamps = deque()

def check_and_wait_for_rate_limit(tokens_used: int):
    current_time = time.time()

    while requests_timestamps and (current_time - requests_timestamps[0] > WINDOW_SECONDS):
        requests_timestamps.popleft()

    while tokens_timestamps and (current_time - tokens_timestamps[0][0] > WINDOW_SECONDS):
        tokens_timestamps.popleft()

    current_requests = len(requests_timestamps)
    current_tokens = sum(t[1] for t in tokens_timestamps)

    if current_requests + 1 > MAX_REQUESTS_PER_MINUTE:
        oldest_request_time = requests_timestamps[0]
        wait_time = WINDOW_SECONDS - (current_time - oldest_request_time)
        if wait_time > 0:
            time.sleep(wait_time)
        return check_and_wait_for_rate_limit(tokens_used)

    if current_tokens + tokens_used > MAX_TOKENS_PER_MINUTE:
        oldest_token_time = tokens_timestamps[0][0]
        wait_time = WINDOW_SECONDS - (current_time - oldest_token_time)
        if wait_time > 0:
            time.sleep(wait_time)
        return check_and_wait_for_rate_limit(tokens_used)

    requests_timestamps.append(current_time)
    tokens_timestamps.append((current_time, tokens_used))

########################################################################
# 2) HÀM ĐỌC TỪ VỰNG CHUYÊN NGÀNH TỪ EXCEL (2 cột: English, Vietnamese)
########################################################################
def load_specialized_dict_from_excel(excel_file):
    if excel_file is None:
        return {}
    df = pd.read_excel(excel_file)
    english_col = "English"
    vietnamese_col = "Vietnamese"

    specialized_dict = {}
    for _, row in df.iterrows():
        eng = str(row[english_col]).strip()
        vie = str(row[vietnamese_col]).strip()
        if eng and vie:
            specialized_dict[eng] = vie
    return specialized_dict

########################################################################
# 3) GỌI OPENAI ĐỂ DỊCH
########################################################################
def translate_text_with_chatgpt(original_text, api_key, global_dict=None):
    partial_dict = {}
    text_lower = original_text.lower()

    if global_dict:
        for eng_term, vie_term in global_dict.items():
            if eng_term.lower() in text_lower:
                partial_dict[eng_term] = vie_term

    system_prompt = (
        "Bạn là một trợ lý AI dịch thuật. Hãy dịch văn bản sau từ tiếng Anh sang tiếng Việt, "
        "ưu tiên dùng đúng các thuật ngữ chuyên ngành (nếu có)."
    )

    dict_prompt = ""
    if partial_dict:
        dict_lines = [f"{k}: {v}" for k, v in partial_dict.items()]
        dict_prompt = (
            "Các thuật ngữ chuyên ngành cần ưu tiên:\n"
            + "\n".join(dict_lines)
            + "\n\n"
        )

    user_prompt = f"{dict_prompt}Văn bản cần dịch:\n{original_text}"

    try:
        client = openai.OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.2,
            max_tokens=2048
        )

        translated_text = response.choices[0].message.content
        total_tokens_used = response.usage.total_tokens if response.usage else 0

        check_and_wait_for_rate_limit(total_tokens_used)
        return translated_text

    except Exception as e:
        st.error(f"Lỗi khi khởi tạo client OpenAI hoặc gọi API: {str(e)}")
        raise


########################################################################
# 4) XỬ LÝ ĐỊNH DẠNG DOCX / PPTX / PDF
########################################################################
def translate_docx(docx_file, translated_texts):
    doc = Document(docx_file)
    all_paragraphs = [p for p in doc.paragraphs if p.text.strip()]

    for i, paragraph in enumerate(all_paragraphs):
        if i < len(translated_texts):
            for run in paragraph.runs:
                run.text = translated_texts[i]

    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output

def translate_pptx(pptx_file, translated_texts):
    pr = Presentation(pptx_file)
    text_shapes = []

    for slide in pr.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text_shapes.append(para)

    for i, para in enumerate(text_shapes):
        if i < len(translated_texts):
            for run in para.runs:
                run.text = translated_texts[i]

    output = BytesIO()
    pr.save(output)
    output.seek(0)
    return output

def translate_pdf(pdf_file, translated_texts):
    reader = PdfReader(pdf_file)
    writer = PdfWriter()

    for i, page in enumerate(reader.pages):
        text = translated_texts[i] if i < len(translated_texts) else ""
        writer.add_page(page)
        writer.pages[i].merge_page(page)
        writer.pages[i].compress_content_streams()

    output = BytesIO()
    writer.write(output)
    output.seek(0)
    return output

########################################################################
# 5) STREAMLIT APP
########################################################################
st.set_page_config(page_title="Auto Translator App with Format Preservation")
st.title("Tự động dịch tài liệu (PPTX, DOCX, PDF) + Giữ định dạng")

api_key = st.text_input("Nhập OpenAI API key của bạn:", type="password")

uploaded_excel_dict = st.file_uploader("Tải lên file Excel chứa thuật ngữ chuyên ngành", type=["xlsx"])
specialized_dict = load_specialized_dict_from_excel(uploaded_excel_dict)

uploaded_file = st.file_uploader("Tải lên file cần dịch (PPTX, DOCX, PDF)", type=["pptx", "docx", "pdf"])

if uploaded_file and api_key:
    file_name = uploaded_file.name
    ext = file_name.split(".")[-1].lower()

    if st.button("Bắt đầu dịch"):
        if ext == "pptx":
            pr = Presentation(uploaded_file)
            all_texts = [para.text for slide in pr.slides for shape in slide.shapes if shape.has_text_frame for para in shape.text_frame.paragraphs]
        elif ext == "docx":
            doc = Document(uploaded_file)
            all_texts = [p.text for p in doc.paragraphs if p.text.strip()]
        elif ext == "pdf":
            reader = PdfReader(uploaded_file)
            all_texts = [page.extract_text() for page in reader.pages if page.extract_text()]
        else:
            st.error("Định dạng không được hỗ trợ.")
            st.stop()

        progress_bar = st.progress(0)
        translated_texts = []
        total = len(all_texts)

        for i, text_chunk in enumerate(all_texts):
            translated = translate_text_with_chatgpt(
                original_text=text_chunk,
                api_key=api_key,
                global_dict=specialized_dict
            )
            translated_texts.append(translated)
            progress_bar.progress((i + 1) / total)

        if ext == "pptx":
            output = translate_pptx(uploaded_file, translated_texts)
            st.download_button("Tải về file PPTX đã dịch", output, "translated.pptx")
        elif ext == "docx":
            output = translate_docx(uploaded_file, translated_texts)
            st.download_button("Tải về file DOCX đã dịch", output, "translated.docx")
        elif ext == "pdf":
            output = translate_pdf(uploaded_file, translated_texts)
            st.download_button("Tải về file PDF đã dịch", output, "translated.pdf")
